#!/usr/bin/env python3
import os
import json
import time
import hashlib
import logging
import sys
from datetime import datetime
from typing import Dict, List, Optional, Any
import traceback
from pathlib import Path
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type

# 로컬 모듈 임포트
script_dir = Path(os.path.dirname(os.path.abspath(__file__)))
sys.path.append(str(script_dir))
try:
    from load_env import load_environment
except ImportError:
    print("load_env.py를 찾을 수 없습니다. 환경 변수를 직접 설정해야 합니다.")

# 환경 변수 로드
try:
    from dotenv import load_dotenv
    if not load_environment():
        print("환경 변수 설정에 실패했습니다. 프로그램을 종료합니다.")
        sys.exit(1)
except ImportError:
    print("python-dotenv 패키지가 설치되어 있지 않습니다. pip install python-dotenv 명령으로 설치하세요.")
    # 환경 변수가 이미 설정되어 있을 수 있으므로 계속 진행

import pinecone
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import io
import google.generativeai as genai
from langchain_google_genai import GoogleGenerativeAIEmbeddings

# 기본 경로 설정
BASE_DIR = script_dir.parent
DATA_DIR = BASE_DIR / "data"
os.makedirs(DATA_DIR, exist_ok=True)

# 로그 설정
log_file = DATA_DIR / "drive_monitor.log"
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("DriveMonitor")

# 상수
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
PINECONE_API_KEY = os.environ.get("PINECONE_API_KEY")
PINECONE_ENVIRONMENT = os.environ.get("PINECONE_ENVIRONMENT")
PINECONE_INDEX = os.environ.get("PINECONE_INDEX")
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
CHECK_INTERVAL = int(os.environ.get("CHECK_INTERVAL", "300"))  # 기본값 5분
DRIVE_FOLDER_ID = os.environ.get("DRIVE_FOLDER_ID", "")  # 모니터링할 폴더 ID
TOKEN_FILE = str(DATA_DIR / 'token.json')
CREDENTIALS_FILE = str(DATA_DIR / 'credentials.json')
CHECKPOINT_FILE = str(DATA_DIR / 'drive_checkpoint.json')
CHUNK_SIZE = 4000  # 모델 요구사항에 맞게 조정

# Gemini 구성
if not GEMINI_API_KEY:
    logger.error("GOOGLE_API_KEY가 설정되지 않았습니다. 프로그램을 종료합니다.")
    sys.exit(1)

genai.configure(api_key=GEMINI_API_KEY)

class DriveMonitor:
    def __init__(self, folder_id=None):
        logger.info("DriveMonitor 초기화 중...")
        
        # 모니터링할 폴더 ID 설정
        self.folder_id = folder_id or DRIVE_FOLDER_ID
        if self.folder_id:
            logger.info(f"특정 폴더(ID: {self.folder_id})만 모니터링합니다.")
        else:
            logger.info("전체 드라이브를 모니터링합니다.")
        
        # 필수 환경 변수 확인
        if not GEMINI_API_KEY:
            logger.error("GEMINI_API_KEY는 필수 환경 변수입니다.")
            sys.exit(1)
        
        # 인증 정보 파일 확인
        if not os.path.exists(CREDENTIALS_FILE):
            logger.error(f"credentials.json 파일이 {CREDENTIALS_FILE}에 없습니다.")
            logger.error("Google Cloud Console에서 OAuth 인증 정보를 다운로드하여 저장하세요.")
            sys.exit(1)
        
        self.credentials = self._get_credentials()
        self.drive_service = build('drive', 'v3', credentials=self.credentials)
        
        # Pinecone 초기화
        logger.info("Pinecone 초기화 중...")
        self.pinecone_client = self._initialize_pinecone()
        self.index = self.pinecone_client.Index(PINECONE_INDEX)
        logger.info(f"Pinecone 인덱스 '{PINECONE_INDEX}' 연결 완료")
        
        # Gemini 임베딩 모델 초기화
        self.embeddings = GoogleGenerativeAIEmbeddings(
            model="models/embedding-001", google_api_key=GEMINI_API_KEY
        )
        self.llm = genai.GenerativeModel('gemini-1.5-pro')
        self.checkpoint = self._load_checkpoint()
        logger.info("DriveMonitor 초기화 완료")
    
    def _get_credentials(self):
        """Get and refresh Google API credentials."""
        creds = None
        if os.path.exists(TOKEN_FILE):
            creds = Credentials.from_authorized_user_info(
                json.load(open(TOKEN_FILE))
            )
        
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                flow = InstalledAppFlow.from_client_secrets_file(
                    CREDENTIALS_FILE, SCOPES
                )
                creds = flow.run_local_server(port=0)
            
            with open(TOKEN_FILE, 'w') as token:
                token.write(creds.to_json())
        
        return creds
    
    def _initialize_pinecone(self):
        """Initialize Pinecone client."""
        global PINECONE_INDEX
        
        if not PINECONE_API_KEY or not PINECONE_ENVIRONMENT:
            logger.error("PINECONE_API_KEY와 PINECONE_ENVIRONMENT는 필수 환경 변수입니다.")
            sys.exit(1)
            
        # Pinecone 클라이언트 초기화
        pc = pinecone.Pinecone(api_key=PINECONE_API_KEY, environment=PINECONE_ENVIRONMENT)
        
        # 인덱스 존재 여부 확인 및 생성
        index_name = PINECONE_INDEX or "kb-profile-data"  # 기본 인덱스 이름
        
        try:
            # 기존 인덱스 리스트 확인
            indexes = pc.list_indexes()
            
            # 인덱스가 없으면 생성
            if index_name not in indexes.names():
                logger.info(f"Pinecone 인덱스 '{index_name}'가 존재하지 않습니다. 생성을 시도합니다.")
                
                # 인덱스 생성
                pc.create_index(
                    name=index_name,
                    dimension=768,  # Gemini embeddings 차원 (기본값)
                    metric="cosine",
                    spec={"serverless": {"cloud": "aws", "region": "us-east-1"}}
                )  
                logger.info(f"Pinecone 인덱스 '{index_name}' 생성 완료")
                
                # 인덱스 준비 대기
                import time
                time.sleep(60)  # 인덱스 초기화 시간 대기
            else:
                logger.info(f"Pinecone 인덱스 '{index_name}'가 이미 존재합니다.")
                
        except Exception as e:
            logger.error(f"Pinecone 인덱스 초기화 중 오류: {str(e)}")
            logger.error(traceback.format_exc())
            sys.exit(1)
        
        # 환경 변수에 인덱스 이름 저장
        os.environ["PINECONE_INDEX"] = index_name
        PINECONE_INDEX = index_name
            
        return pc
    
    def _load_checkpoint(self) -> Dict[str, str]:
        """Load the last checkpoint for file modification times."""
        if os.path.exists(CHECKPOINT_FILE):
            with open(CHECKPOINT_FILE, 'r') as f:
                return json.load(f)
        return {}
    
    def _save_checkpoint(self):
        """Save current checkpoint to file."""
        with open(CHECKPOINT_FILE, 'w') as f:
            json.dump(self.checkpoint, f)
    
    def get_file_changes(self) -> List[Dict[str, Any]]:
        """Get files that have been modified since last check."""
        query = "trashed=false and mimeType!='application/vnd.google-apps.folder'"
        
        # 특정 폴더만 모니터링하는 경우
        if self.folder_id:
            query += f" and '{self.folder_id}' in parents"
            
        results = self.drive_service.files().list(
            q=query,
            fields="files(id, name, mimeType, modifiedTime, md5Checksum)"
        ).execute()
        
        files = results.get('files', [])
        changed_files = []
        
        for file in files:
            file_id = file['id']
            modified_time = file['modifiedTime']
            checksum = file.get('md5Checksum', '')
            
            # Check if file is new or modified
            if (file_id not in self.checkpoint or 
                self.checkpoint[file_id] != checksum):
                changed_files.append(file)
                self.checkpoint[file_id] = checksum
        
        return changed_files
    
    def download_file(self, file_id: str, mime_type: str = None) -> io.BytesIO:
        """Download a file from Google Drive.
        
        If mime_type is provided, attempts to export Google Workspace files to that format.
        """
        file_content = io.BytesIO()
        
        try:
            # For Google Workspace formats, use export
            if mime_type and 'application/vnd.google-apps' in mime_type:
                export_mime_type = self._get_export_mime_type(mime_type)
                request = self.drive_service.files().export_media(fileId=file_id, mimeType=export_mime_type)
            else:
                # For regular files, use get_media
                request = self.drive_service.files().get_media(fileId=file_id)
            
            downloader = MediaIoBaseDownload(file_content, request)
            
            done = False
            while not done:
                _, done = downloader.next_chunk()
            
            file_content.seek(0)
        except Exception as e:
            print(f"Error downloading file {file_id}: {str(e)}")
            # Return empty content on error
            file_content = io.BytesIO()
        
        return file_content
    
    def _get_export_mime_type(self, google_mime_type: str) -> str:
        """Get the appropriate export MIME type for Google Workspace files."""
        export_formats = {
            'application/vnd.google-apps.document': 'text/plain',
            'application/vnd.google-apps.spreadsheet': 'text/csv',
            'application/vnd.google-apps.presentation': 'text/plain',
            'application/vnd.google-apps.drawing': 'image/png',
            'application/vnd.google-apps.script': 'application/vnd.google-apps.script+json'
        }
        
        return export_formats.get(google_mime_type, 'text/plain')
    
    def extract_text(self, file_content: io.BytesIO, mime_type: str, file_name: str) -> str:
        """Extract text from a file based on its MIME type."""
        # Save the content to a temporary file to process it
        import tempfile
        import os
        
        temp_dir = tempfile.gettempdir()
        temp_file_path = os.path.join(temp_dir, file_name)
        
        with open(temp_file_path, 'wb') as temp_file:
            temp_file.write(file_content.getvalue())
        
        extracted_text = ""
        
        try:
            # Process PDF files
            if 'pdf' in mime_type or file_name.lower().endswith('.pdf'):
                try:
                    import PyPDF2
                    with open(temp_file_path, 'rb') as pdf_file:
                        pdf_reader = PyPDF2.PdfReader(pdf_file)
                        for page_num in range(len(pdf_reader.pages)):
                            page = pdf_reader.pages[page_num]
                            extracted_text += page.extract_text() + "\n\n"
                except ImportError:
                    extracted_text = f"Error: PyPDF2 library not found. Please install with: pip install PyPDF2"
            
            # Process Word documents
            elif 'word' in mime_type or file_name.lower().endswith(('.docx', '.doc')):
                try:
                    import docx
                    doc = docx.Document(temp_file_path)
                    extracted_text = "\n".join([para.text for para in doc.paragraphs])
                except ImportError:
                    extracted_text = f"Error: python-docx library not found. Please install with: pip install python-docx"
            
            # Process Excel files
            elif 'excel' in mime_type or 'spreadsheet' in mime_type or file_name.lower().endswith(('.xlsx', '.xls')):
                try:
                    import pandas as pd
                    if file_name.lower().endswith('.xlsx'):
                        excel_data = pd.read_excel(temp_file_path, engine='openpyxl')
                    else:
                        excel_data = pd.read_excel(temp_file_path)
                    extracted_text = excel_data.to_string()
                except ImportError:
                    extracted_text = f"Error: pandas or openpyxl libraries not found. Please install with: pip install pandas openpyxl"
            
            # Process PowerPoint files
            elif 'presentation' in mime_type or file_name.lower().endswith(('.pptx', '.ppt')):
                try:
                    import pptx
                    presentation = pptx.Presentation(temp_file_path)
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted_text += shape.text + "\n"
                        extracted_text += "\n"
                except ImportError:
                    extracted_text = f"Error: python-pptx library not found. Please install with: pip install python-pptx"
            
            # Process plain text files
            elif 'text' in mime_type or file_name.lower().endswith(('.txt', '.md', '.csv')):
                with open(temp_file_path, 'r', encoding='utf-8', errors='ignore') as text_file:
                    extracted_text = text_file.read()
            
            # Google Docs format
            elif 'application/vnd.google-apps.document' in mime_type:
                # Need to export Google Docs to a readable format first
                # This is just a placeholder - the actual implementation would use the Drive API export method
                extracted_text = f"Google Docs file: {file_name}. Use Drive API export to convert to text."
            
            # Default handling for unknown file types
            else:
                extracted_text = f"Unsupported file type: {mime_type} for file {file_name}"
        
        except Exception as e:
            extracted_text = f"Error extracting text from {file_name}: {str(e)}"
        
        # Clean up the temporary file
        try:
            os.remove(temp_file_path)
        except:
            pass
        
        return extracted_text
    
    def chunk_text(self, text: str) -> List[str]:
        """Break text into manageable chunks for embedding."""
        chunks = []
        current_length = 0
        current_chunk = []
        
        for line in text.split('\n'):
            line_length = len(line)
            
            if current_length + line_length > CHUNK_SIZE:
                chunks.append('\n'.join(current_chunk))
                current_chunk = [line]
                current_length = line_length
            else:
                current_chunk.append(line)
                current_length += line_length
        
        if current_chunk:
            chunks.append('\n'.join(current_chunk))
        
        return chunks
    
    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=4, max=60),
        retry=retry_if_exception_type((IOError, ConnectionError))
    )
    def generate_summary(self, text: str) -> str:
        """Generate a summary using Gemini with retry logic."""
        try:
            logger.debug(f"Generating summary for text of length {len(text)}")
            response = self.llm.generate_content(
                f"""Please provide a concise summary of the following content:
                
                {text[:8000]}  # Limit text length for summary generation
                
                Summary:"""
            )
            return response.text
        except Exception as e:
            logger.error(f"Error generating summary: {str(e)}")
            return f"Error generating summary: {str(e)}"
    
    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=4, max=60),
        retry=retry_if_exception_type((IOError, ConnectionError))
    )
    def extract_keywords(self, text: str) -> List[str]:
        """Extract keywords using Gemini with retry logic."""
        try:
            logger.debug(f"Extracting keywords from text of length {len(text)}")
            response = self.llm.generate_content(
                f"""Extract the 5-10 most important keywords from this content. 
                Return them as a comma-separated list with no explanations:
                
                {text[:8000]}  # Limit text length for keyword extraction
                
                Keywords:"""
            )
            keywords = response.text.strip().split(',')
            return [k.strip() for k in keywords]
        except Exception as e:
            logger.error(f"Error extracting keywords: {str(e)}")
            return ["error_extracting_keywords"]
    
    def infer_doc_type(self, file_name: str, text: str) -> str:
        """Infer document type based on content and file name."""
        file_name_lower = file_name.lower()
        text_sample = text[:1000].lower()  # Use a sample of the text
        
        if any(term in file_name_lower for term in ['resume', 'cv']):
            return 'resume'
        elif any(term in file_name_lower for term in ['pub', 'paper', 'article']):
            return 'publication'
        elif any(term in file_name_lower for term in ['profile', 'bio']):
            return 'profile'
        
        # If filename doesn't give clear indication, try content analysis
        if any(term in text_sample for term in ['experience', 'education', 'skills']):
            return 'resume'
        elif any(term in text_sample for term in ['abstract', 'introduction', 'methodology']):
            return 'publication'
        
        # Default doctype
        return 'document'
    
    def create_vector_record(self, chunk: str, doc_type: str, 
                           summary: str, keywords: List[str], 
                           source: str) -> Dict[str, Any]:
        """Create the vector record with the provided structure."""
        return {
            "text": chunk,
            "doc_type": doc_type,
            "summary": summary,
            "keywords": keywords,
            "source": source
        }
    
    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=4, max=60),
        retry=retry_if_exception_type((IOError, ConnectionError))
    )
    def process_file(self, file: Dict[str, Any]):
        """Process a single file from Google Drive with retry logic."""
        file_id = file['id']
        file_name = file['name']
        mime_type = file['mimeType']
        
        logger.info(f"Processing file: {file_name} (ID: {file_id})")
        
        try:
            # Download file (passing mime_type for Google Workspace formats)
            file_content = self.download_file(file_id, mime_type)
            
            if file_content.getbuffer().nbytes == 0:
                logger.warning(f"Empty content downloaded for file: {file_name}. Skipping.")
                return
            
            # Extract text
            text = self.extract_text(file_content, mime_type, file_name)
            
            if not text or len(text) < 10:
                logger.warning(f"Insufficient text extracted from file: {file_name}. Skipping.")
                return
            
            # Identify document type
            doc_type = self.infer_doc_type(file_name, text)
            logger.info(f"Identified document type for {file_name}: {doc_type}")
            
            # Generate summary for the entire document
            logger.info(f"Generating summary for {file_name}")
            summary = self.generate_summary(text)
            
            # Extract keywords for the entire document
            logger.info(f"Extracting keywords for {file_name}")
            keywords = self.extract_keywords(text)
            
            # Chunk text for embedding
            logger.info(f"Chunking text for {file_name}")
            chunks = self.chunk_text(text)
            logger.info(f"Created {len(chunks)} chunks for {file_name}")
            
            # Process each chunk
            for i, chunk in enumerate(chunks):
                try:
                    # Create unique ID for this chunk
                    chunk_id = f"{file_id}-chunk-{i}"
                    
                    # Create embedding
                    logger.debug(f"Creating embedding for chunk {i+1}/{len(chunks)} of {file_name}")
                    embedding = self.embeddings.embed_query(chunk)
                    
                    # Create record
                    record = self.create_vector_record(
                        chunk=chunk,
                        doc_type=doc_type,
                        summary=summary,
                        keywords=keywords,
                        source=file_name
                    )
                    
                    # Upsert to Pinecone
                    logger.debug(f"Upserting chunk {i+1}/{len(chunks)} for {file_name} to Pinecone")
                    self.index.upsert(
                        vectors=[(chunk_id, embedding, record)]
                    )
                    
                    logger.info(f"Processed chunk {i+1}/{len(chunks)} for {file_name}")
                except Exception as e:
                    logger.error(f"Error processing chunk {i+1}/{len(chunks)} for {file_name}: {str(e)}")
                    logger.error(traceback.format_exc())
        except Exception as e:
            logger.error(f"Error processing file {file_name}: {str(e)}")
            logger.error(traceback.format_exc())
            raise  # Re-raise for retry mechanism
    
    def monitor_and_process(self, check_interval: int = None):
        """Main loop to monitor Drive for changes and process files."""
        if check_interval is None:
            check_interval = CHECK_INTERVAL
            
        logger.info(f"구글 드라이브 모니터링 시작 (확인 간격: {check_interval}초)")
        
        try:
            while True:
                logger.info(f"변경사항 확인 중: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
                try:
                    changed_files = self.get_file_changes()
                    
                    if changed_files:
                        logger.info(f"{len(changed_files)}개의 변경된 파일 발견")
                        for file in changed_files:
                            try:
                                self.process_file(file)
                            except Exception as e:
                                logger.error(f"재시도 후에도 파일 처리 실패: {file['name']}: {str(e)}")
                    else:
                        logger.info("변경사항 없음")
                    
                    # Save checkpoint after processing
                    self._save_checkpoint()
                    
                except Exception as e:
                    logger.error(f"변경 감지 사이클 중 오류 발생: {str(e)}")
                    logger.error(traceback.format_exc())
                
                logger.info(f"{check_interval}초 후 다시 확인")
                time.sleep(check_interval)
                
        except KeyboardInterrupt:
            logger.info("사용자에 의해 모니터링 중단")
            self._save_checkpoint()
        except Exception as e:
            logger.error(f"모니터링 루프에서 예상치 못한 오류 발생: {str(e)}")
            logger.error(traceback.format_exc())
            self._save_checkpoint()

if __name__ == "__main__":
    try:
        logger.info("드라이브 모니터 애플리케이션 시작")
        # 명령줄 인자로 폴더 ID를 받을 수 있도록 함
        import argparse
        parser = argparse.ArgumentParser(description='Google Drive 변경사항 모니터링 및 벡터 생성')
        parser.add_argument('--folder', '-f', help='모니터링할 Google Drive 폴더 ID')
        args = parser.parse_args()
        
        monitor = DriveMonitor(folder_id=args.folder)
        monitor.monitor_and_process()
    except Exception as e:
        logger.critical(f"치명적 오류: {str(e)}")
        logger.critical(traceback.format_exc())
